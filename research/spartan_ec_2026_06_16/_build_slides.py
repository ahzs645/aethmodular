"""Builds spartan_ec_weekly_2026_06_16.pptx — this week's simple update deck.

Graph-per-slide + short bullets + the confirmation points. Consumes the PNGs already
rendered by the 02/03/04 notebooks into figures/. Style mirrors
research/ftir_hips_chem/scripts/build_warren_meeting_deck.py (python-pptx, 16:9).
Run after the notebooks have produced their figures.
"""
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
FIG = HERE / "figures"
DECK = HERE / "spartan_ec_weekly_2026_06_16.pptx"

SLIDE_W, SLIDE_H = 13.333, 7.5
TITLE_PT, BULLET_PT, CAPTION_PT = 30, 20, 13


def title_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    box = s.shapes.add_textbox(Inches(0.6), Inches(2.6), Inches(12.1), Inches(1.6))
    tf = box.text_frame; tf.word_wrap = True; tf.text = title
    tf.paragraphs[0].runs[0].font.size = Pt(40)
    tf.paragraphs[0].runs[0].font.bold = True
    box2 = s.shapes.add_textbox(Inches(0.6), Inches(4.2), Inches(12.1), Inches(1.4))
    tf2 = box2.text_frame; tf2.word_wrap = True; tf2.text = subtitle
    tf2.paragraphs[0].runs[0].font.size = Pt(20)
    return s


def add_title(s, text, size=TITLE_PT):
    box = s.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.85))
    tf = box.text_frame; tf.word_wrap = True; tf.text = text
    for r in tf.paragraphs[0].runs:
        r.font.size = Pt(size); r.font.bold = True


def add_bullets(s, items, *, left=0.6, top=1.2, width=12.1, height=5.6, size=BULLET_PT):
    box = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if it.startswith("\t"):
            p.text = "      – " + it.lstrip("\t "); p.level = 1
        else:
            p.text = "•  " + it
        for r in p.runs:
            r.font.size = Pt(size)
        p.space_after = Pt(6)


def add_image(s, png, *, left=0.4, top=1.15, width=12.5, height=4.9):
    p = Path(png)
    with Image.open(p) as im:
        sw, sh = im.size
    sr, br = sw / sh, width / height
    if sr >= br:
        dw, dh = width, width / sr
    else:
        dh, dw = height, height * sr
    s.shapes.add_picture(str(p), Inches(left + (width - dw) / 2),
                         Inches(top + (height - dh) / 2), Inches(dw), Inches(dh))


def add_caption(s, text, *, top=6.2):
    box = s.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(1.1))
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + line
        for r in p.runs:
            r.font.size = Pt(CAPTION_PT)


prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)

# 1 — title
title_slide(prs, "SPARTAN EC — Weekly Update",
            "Ahmad Jalil · 2026-06-16 · toward a SPARTAN-network EC paper (AAAR abstract due Jul 24)")

# 2 — confirmation points
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "This week — confirmation points")
add_bullets(s, [
    "Confirmed the 2003 HIPS cutoff straight from Warren White (2016): “stable calibration since 2003.”",
    "Confirmed SPARTAN EC = FTIR. There is NO thermal/TOR EC in SPARTAN (no quartz filters).",
    "\tThe “thermal” stars on the EC/OC plot were optical BC (HIPS/SSR) — dropped them.",
    "Reproduced fAbs/EC by site exactly: Beijing 10.3, Pasadena 9.4, Delhi 8.5.",
    "Corrected lot for Addis: 251 (184 samples), not 256.",
    "Got the 5 Adama AMOD quartz-filter TOR samples — our first real thermal EC near Addis.",
])

# 3 — SPARTAN carbon = FTIR (the method confirmation)
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "SPARTAN carbon is FTIR — no thermal/TOR EC")
add_bullets(s, [
    "EC PM2.5 and OC PM2.5  →  methods 217 / 218  =  FTIR (the functional-group method under review)",
    "BC PM2.5  →  methods 219 / 221 = HIPS (optical); 220 = HIPS–SSR curve",
    "Equivalent BC PM2.5  →  methods 212/214/215/216 = Smoke Stain Reflectometer",
    "Verified across all 4 focus sites (ETAD, CHTS, INDH, USPA).",
    "So the only independent check on FTIR-EC is the optical HIPS — which is why the fAbs/tau cross-plots carry the argument.",
])

# 4 — cross-plot fixed
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Volume-free cross-plot — reference lines fixed")
add_image(s, FIG / "fig04A_tau_vs_ecmass.png")
add_caption(s, [
    "Dropped the “1:1” line — meaningless (tau is dimensionless, EC-mass is µg). tau = fAbs·V/A confirmed unitless.",
    "Gray = the real IMPROVE through-origin slope; black dashed = MAC=10 reference. Addis sits on a steeper line.",
])

# 5 — fAbs/EC by site
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "fAbs/EC by site (effective MAC)")
add_image(s, FIG / "fig04B_fabs_ec_by_site.png")
add_caption(s, [
    "Reproduces the meeting values (Beijing 10.3, Pasadena 9.4, Delhi 8.5). Addis ≈ 10 — not special on this metric.",
    "IMPROVE median ≈ 12. Caveat: if FTIR-EC is biased low, this “MAC” is biased high.",
])

# 6 — EC/OC by site
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "EC/OC by site (FTIR EC only)")
add_image(s, FIG / "fig04C_ec_oc_by_site.png")
add_caption(s, [
    "FTIR EC only (the “thermal” stars are gone — they weren’t thermal). Addis EC/OC highest of the four.",
    "IMPROVE much lower — the compositional difference behind the whole investigation.",
])

# 7 — fractional EC
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Fractional EC — focus sites vs IMPROVE")
add_image(s, FIG / "fig04D_fractional_ec.png")
add_caption(s, [
    "IMPROVE fractionally much lower than every SPARTAN focus site; Addis the most extreme.",
    "All-sites version + ETBI need the broader public chemspec — ETBI has HIPS but no EC yet (Chris/Oxford upload).",
])

# 8 — biomass calibration (text, plan + baseline)
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Biomass-burning calibration test (top priority, with Mona)")
add_bullets(s, [
    "Plan: build a biomass-only calibration on lot 251 (Chelsea’s guidance), apply to Addis, compare EC vs the general calibration.",
    "Tests our original EC-paper point: smoke samples should be calibrated with smoke samples.",
    "Baseline is built: lot-251 Addis fit  HIPS(BC-eq) = 0.43·EC + 2.65  (r = 0.89).",
    "\tThat +2.65 µg/m³ ≈ 26.5 Mm⁻¹ intercept matches the ~28 Mm⁻¹ Deming intercept — the gap is in the intercept.",
    "Status: awaiting the biomass-only calibration + per-filter EC from the FTIR tool (Mona).",
])

# 9 — Adama TOR char/soot
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Adama TOR (5 quartz filters) — Han char vs soot")
add_image(s, FIG / "fig03_adama_char_soot.png")
add_caption(s, [
    "char-EC/soot-EC is low (0.02–0.58) → the EC present is soot-leaning. But OC/EC is high (4.6–7.2) + large pyrolysis → organic-rich / biomass-influenced.",
    "Two metrics disagree — that’s the finding: exactly the regime where the FTIR calibration is expected to struggle.",
])

# 10 — next steps
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Next steps")
add_bullets(s, [
    "Mona: biomass-only lot-251 calibration + per-filter EC; then Addis spectra to compare.",
    "Predict EC for the 5 Adama filters with general vs wood-smoke calibration; compare to TOR EC.",
    "Add the Han (char/soot) and KBr-pellet papers to Zotero so we can cite them.",
    "Charcoal KBr prep (~0.4–0.5 mg / 300 mg KBr); dry separately in the muffle furnace.",
    "AAAR abstract (due Jul 24): frame around SPARTAN EC.",
])

prs.save(str(DECK))
print("wrote", DECK.name, "—", len(prs.slides._sldIdLst), "slides")
