"""Build the Warren-meeting deck following Ann's revised (purple-comment) feedback.

Slide order (Ann's email, May 2026 revision):
  1. Intro: dataset + problem
  2. Data and tools used
  3. Fabs (1/Mm) vs FTIR EC (ug/m3), four SPARTAN sites
  4. Fabs (1/Mm) vs FTIR EC mass on filter (ug), four sites
  5. Fabs (1/Mm) vs Aethalometer IR BCc
  6. Iron: Fe vs Fabs + Fe/EC distribution (with IMPROVE shading)
  7. Seasonality + checklist of ruled-out variables
  8. Four SPARTAN sites overlaid on IMPROVE (EC mass on filter)
  9. Per-site, 5-95 percentile shading on both axes
 10. SPARTAN R/T calibration space, split by site
 10b. All four sites overlaid: Addis transmittance offset
 11. Questions for Warren

Style: white background, large fonts (Warren-friendly per Ann's note that he
cares about data, not beauty).
"""
from __future__ import annotations
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

import os
# Resolve paths so this script works whether run from the host or the sandbox.
_HOST = Path("/Users/ahmadjalil/github/aethmodular")
_SANDBOX = Path("/sessions/gallant-zealous-franklin/mnt/research").parent / "research"
# Prefer whichever exists (sandbox sees research/ at /sessions/.../mnt/research)
if _HOST.exists():
    REPO_ROOT = _HOST
elif Path("/sessions/gallant-zealous-franklin/mnt/research").exists():
    # /sessions/.../mnt/research IS the research/ dir on host
    REPO_ROOT = Path("/sessions/gallant-zealous-franklin/mnt")
else:
    REPO_ROOT = _HOST  # fallback for error message
RESEARCH = REPO_ROOT / "research" / "ftir_hips_chem"
FIG_DIR = RESEARCH / "output" / "warren_meeting" / "figures"
DECK_PATH = RESEARCH / "output" / "warren_meeting" / "warren_meeting_for_warren.pptx"
# Cache of caption-stripped McDade screenshots (re-generated on first build).
OUTPUT_DIR_CROPPED = RESEARCH / "output" / "warren_meeting" / "figures" / "_mcdade_cropped"

# Filter-pattern source images (Ann's "add a slide showing SPARTAN filters" email).
# SPARTAN photos sent by Lindsay Kline (UC Davis AQRC) on 4/30. Mounted into the
# sandbox at /sessions/.../mnt/filter when working from bash.
_FILTER_HOST = Path("/Users/ahmadjalil/Downloads/filter")
_FILTER_SANDBOX = Path("/sessions/gallant-zealous-franklin/mnt/filter")
FILTER_DIR = _FILTER_SANDBOX if _FILTER_SANDBOX.exists() else _FILTER_HOST
SPARTAN_FILTER_PHOTOS = [
    FILTER_DIR / "image004.jpg",
    FILTER_DIR / "image005.jpg",
    FILTER_DIR / "image006.jpg",
]
# McDade et al. (2009) IMPROVE filter pattern figures.
# Source: ~/Downloads/Improve Filters/ (PDF + screenshots taken from
# Particulate_Matter_Sample_Deposit_Geometry_and_Eff.pdf).
# Three sub-figures used here, mirroring the SPARTAN-photo slide layout:
#   Fig 1 — filter holder support screen (25mm)
#   Fig 4 — sampled IMPROVE Teflon filter (full diameter)
#   Fig 2 — close-up of dot deposition pattern (~0.013" per dot)
_MCDADE_HOST = Path("/Users/ahmadjalil/Downloads/Improve Filters")
_MCDADE_SANDBOX = Path("/sessions/gallant-zealous-franklin/mnt/Improve Filters")
MCDADE_DIR = _MCDADE_SANDBOX if _MCDADE_SANDBOX.exists() else _MCDADE_HOST


def _resolve_mcdade_figs():
    """Match McDade screenshots by timestamp substring.

    macOS Screenshot filenames use a narrow no-break space (U+202F) before "PM",
    so glob/exact-match on a regular-space filename fails. Match on the unique
    timestamp instead.
    """
    if not MCDADE_DIR.exists():
        return []
    by_substr = {p.name: p for p in MCDADE_DIR.iterdir() if p.suffix.lower() in {".png", ".jpg", ".jpeg"}}
    def find(stamp):
        for name, path in by_substr.items():
            if stamp in name:
                return path
        return MCDADE_DIR / f"missing-{stamp}.png"
    return [
        (find("11.29.04"),
         "Fig 1 — IMPROVE sampler filter holder screen (25 mm)"),
        (find("11.29.10"),
         "Fig 4 — Sampled IMPROVE Teflon filter (full sample diameter)"),
        (find("11.28.59"),
         "Fig 2 — Sampled IMPROVE filter close-up; each dot ≈ 0.013″"),
    ]


MCDADE_FIGS = _resolve_mcdade_figs()
# Optional single overview PNG. Kept for backward compatibility with the
# original build script — used only if MCDADE_FIGS images are not present.
MCDADE_IMPROVE_FILTER_PNG = REPO_ROOT / "tmp_warren_pages" / "mcdade_improve_filter.png"

# 16:9
W = 13.333
H = 7.5

# Larger sizes per Ann's "current font is too small" note.
TITLE_PT = 34
BODY_PT = 22
SUB_PT = 18
CAPTION_PT = 16

# A single accent color used sparingly (an Addis-amber tone, matches notebook).
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_MUTED = RGBColor(0x55, 0x55, 0x55)


def add_title(slide, text, *, size=TITLE_PT, top=0.25, color=TEXT_DARK):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(W - 1.0), Inches(0.85))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.text = text
    p = tf.paragraphs[0]
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.name = "Calibri"
        run.font.color.rgb = color


def add_accent_underline(slide, top=1.05, width=1.6):
    """Small accent rule under the title (NOT a full-width line — that's the AI tell)."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(top), Inches(width), Inches(0.06)
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT


def crop_caption_strip(src_png, out_png):
    """Auto-crop the 'Figure N. ...' caption strip baked into McDade screenshots.

    Strategy: walk up from the bottom looking for the last row that contains
    near-black pixels (the figure's content). Trim a small margin below it so
    the caption text — which sits in a pure-white band — is dropped. Falls
    back to a fixed 12% bottom crop if detection fails.
    """
    from PIL import Image as _Image
    import numpy as np
    out = Path(out_png)
    if out.exists():
        return out
    out.parent.mkdir(parents=True, exist_ok=True)
    with _Image.open(src_png) as im:
        arr = np.asarray(im.convert("L"))  # grayscale
    h = arr.shape[0]
    # Rows with any "darkish" pixels (content). Threshold tuned for paper figures.
    has_content = (arr < 200).any(axis=1)
    # Walk up from bottom past the trailing all-white margin into the caption,
    # then past the caption (also dark), then past the whitespace gap above it.
    # Heuristic: find the last contiguous content block, drop everything below
    # the gap that precedes it.
    rows = np.where(has_content)[0]
    if len(rows) == 0:
        bottom = int(h * 0.88)
    else:
        # Find gaps in the content rows; the last big gap separates the figure
        # from the caption.
        diffs = np.diff(rows)
        big_gaps = np.where(diffs > max(8, h // 50))[0]
        if len(big_gaps) > 0:
            # rows[big_gaps[-1]] is the last row of content before the gap before
            # the caption. Crop just above the caption (= just below that row).
            bottom = int(rows[big_gaps[-1]] + max(4, h // 200))
        else:
            bottom = int(h * 0.88)
    bottom = max(int(h * 0.55), min(bottom, h))
    with _Image.open(src_png) as im:
        cropped = im.crop((0, 0, im.size[0], bottom))
        cropped.save(out)
    return out


def add_image_fit(slide, png, *, left, top, width, height):
    p = Path(png)
    if not p.exists():
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = f"[FIGURE MISSING: {png}]"
        for r in tf.paragraphs[0].runs:
            r.font.size = Pt(BODY_PT)
            r.font.italic = True
        return
    with Image.open(p) as im:
        sw, sh = im.size
    sr = sw / sh
    br = width / height
    if sr >= br:
        dw = width
        dh = width / sr
    else:
        dh = height
        dw = height * sr
    dl = left + (width - dw) / 2
    dt = top + (height - dh) / 2
    slide.shapes.add_picture(str(p), Inches(dl), Inches(dt), width=Inches(dw), height=Inches(dh))


def add_bullets(slide, items, *, left=0.7, top=1.5, width=W - 1.4, height=5.5,
                size=BODY_PT, line_spacing=1.25):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.1)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        is_sub = item.startswith("\t")
        text = item.lstrip("\t ")
        if is_sub:
            p.text = "    – " + text
            p.level = 1
            sz = size - 4
        else:
            p.text = "• " + text
            sz = size
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        for r in p.runs:
            r.font.size = Pt(sz)
            r.font.name = "Calibri"
            r.font.color.rgb = TEXT_DARK


def add_takeaway(slide, text, *, top=6.55, height=0.85, size=20):
    """Single key-takeaway line pinned at the bottom of a figure slide."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(top), Inches(0.08), Inches(height)
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    box = slide.shapes.add_textbox(
        Inches(0.6), Inches(top), Inches(W - 1.1), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.05)
    tf.text = text
    p = tf.paragraphs[0]
    p.line_spacing = 1.15
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.name = "Calibri"
        r.font.color.rgb = TEXT_DARK


def add_speaker_notes(slide, text):
    """Write `text` into the slide's notes pane (stays out of the visible slide)."""
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(H - 0.35), Inches(W - 0.8), Inches(0.3))
    tf = box.text_frame
    tf.text = text
    for r in tf.paragraphs[0].runs:
        r.font.size = Pt(11)
        r.font.italic = True
        r.font.color.rgb = TEXT_MUTED


# ── Slide 1: intro / cover ────────────────────────────────────────────────
def slide_intro(prs, blank):
    s = prs.slides.add_slide(blank)
    # left accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(H))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    # Title
    title_box = s.shapes.add_textbox(Inches(0.9), Inches(1.3), Inches(W - 1.5), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.text = "Why is HIPS Fabs at Addis Ababa so high?"
    for r in tf.paragraphs[0].runs:
        r.font.size = Pt(46)
        r.font.bold = True
        r.font.name = "Calibri"
        r.font.color.rgb = TEXT_DARK
    # Subtitle
    sub = s.shapes.add_textbox(Inches(0.9), Inches(3.2), Inches(W - 1.5), Inches(0.8))
    tf = sub.text_frame
    tf.text = "Cross-network comparison with IMPROVE"
    for r in tf.paragraphs[0].runs:
        r.font.size = Pt(24)
        r.font.color.rgb = TEXT_MUTED

    # The "problem" framing
    prob = s.shapes.add_textbox(Inches(0.9), Inches(4.4), Inches(W - 1.5), Inches(2.4))
    tf = prob.text_frame
    tf.word_wrap = True
    items = [
        "Four SPARTAN sites: Addis Ababa, Delhi, JPL/Pasadena, Beijing",
        "HIPS Fabs at Addis is anomalously high vs the other three sites and vs IMPROVE",
    ]
    for i, t in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + t
        p.line_spacing = 1.25
        p.space_after = Pt(6)
        for r in p.runs:
            r.font.size = Pt(20)
            r.font.color.rgb = TEXT_DARK

    # Presenter
    pres = s.shapes.add_textbox(Inches(0.9), Inches(6.7), Inches(W - 1.5), Inches(0.7))
    tf = pres.text_frame
    tf.text = "Ahmad Jalil  ·  UNBC / NASA MAIA / SPARTAN  ·  May 2026"
    for r in tf.paragraphs[0].runs:
        r.font.size = Pt(14)
        r.font.color.rgb = TEXT_MUTED


# ── Slide 2: data and tools ───────────────────────────────────────────────
def slide_data_tools(prs, blank):
    s = prs.slides.add_slide(blank)
    add_title(s, "Data and tools used")
    add_accent_underline(s)
    items = [
        "SPARTAN sites (PM2.5 quartz + PTFE filters):",
        "\tAddis Ababa  ·  Delhi  ·  JPL / Pasadena  ·  Beijing",
        "Filter measurements:",
        "\tFTIR EC (mass on filter and µg/m³)",
        "\tHIPS Fabs (raw, kept in 1/Mm, NOT divided by MAC = 10)",
        "\tChemSpec PM2.5 metals (Fe used here)",
        "Continuous: microAeth MA350 IR BCc (1-min, averaged over filter day)",
    ]
    add_bullets(s, items, top=1.4, size=BODY_PT)


# ── Slide 2.5: Addis concentrations in context (Ann's 5/8 email) ──────────
def slide_addis_context(prs, blank):
    s = prs.slides.add_slide(blank)
    add_title(s, "Putting Addis concentrations in context")
    add_accent_underline(s)
    items = [
        "Addis Ababa, Dec 2022 – Aug 2024 (this study, 24-h means):",
        "\tFTIR EC: 5.1 µg/m³   ·   HIPS BC (MAC=10): 4.9 µg/m³   ·   Aethalometer BC: 7.8 µg/m³",
        "Why so high? Ethiopian highland city at 2,300 m:",
        "\tDistributed household charcoal & biomass burning",
        "\tOlder vehicle fleet, limited emissions controls",
        "\tLocal & transported mineral dust",
        "How it compares:",
        "\tAmong highest urban BC measured across African cities (Anand et al. 2024)",
        "\tBC across the global south chronically underestimated in inventories (Ren et al. 2025)",
        "\tWithin SPARTAN: Addis FTIR EC ≫ Beijing, JPL; comparable to Delhi",
    ]
    add_bullets(s, items, top=1.4, size=BODY_PT)


# ── Filter-pattern slides (Ann's 4/30 email) ──────────────────────────────
def slide_spartan_filters(prs, blank):
    s = prs.slides.add_slide(blank)
    add_title(s, "SPARTAN PTFE filters: support-mesh grid pattern")
    add_accent_underline(s)
    photos = [p for p in SPARTAN_FILTER_PHOTOS if Path(p).exists()]
    if photos:
        n = len(photos)
        gap = 0.3
        total_w = W - 1.0
        photo_w = (total_w - gap * (n - 1)) / n
        photo_h = 5.6
        photo_top = 1.35
        for i, photo in enumerate(photos):
            left = 0.5 + i * (photo_w + gap)
            add_image_fit(s, photo, left=left, top=photo_top, width=photo_w, height=photo_h)
    else:
        msg = (
            f"[FIGURE MISSING: copy SPARTAN filter JPGs to {FILTER_DIR} "
            "(image004.jpg, image005.jpg, image006.jpg) and re-run build_deck.py]"
        )
        box = s.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(W - 1.0), Inches(1.5))
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = msg
        for r in tf.paragraphs[0].runs:
            r.font.size = Pt(20)
            r.font.italic = True
            r.font.color.rgb = TEXT_MUTED


def slide_mcdade_improve(prs, blank):
    s = prs.slides.add_slide(blank)
    add_title(s, "IMPROVE filter pattern: McDade et al. (2009)")
    add_accent_underline(s)

    # Preferred path: 3 sub-figures from McDade 2009 side-by-side.
    fig_paths = [(p, c) for p, c in MCDADE_FIGS if Path(p).exists()]
    if fig_paths:
        n = len(fig_paths)
        gap = 0.35
        total_w = W - 1.0
        photo_w = (total_w - gap * (n - 1)) / n
        photo_h = 5.5
        photo_top = 1.35
        cache = OUTPUT_DIR_CROPPED
        for i, (p, _cap) in enumerate(fig_paths):
            left = 0.5 + i * (photo_w + gap)
            cropped = crop_caption_strip(p, cache / Path(p).with_suffix(".cropped.png").name)
            add_image_fit(s, cropped, left=left, top=photo_top, width=photo_w, height=photo_h)
        # Single citation strip pinned at the bottom
        cite = s.shapes.add_textbox(Inches(0.4), Inches(7.18), Inches(W - 0.8), Inches(0.25))
        tf = cite.text_frame
        tf.text = (
            "McDade et al. (2009), Particulate matter sample deposit geometry and effective "
            "filter face velocities."
        )
        for r in tf.paragraphs[0].runs:
            r.font.size = Pt(11)
            r.font.italic = True
            r.font.color.rgb = TEXT_MUTED
        return

    # Backward-compat: single overview PNG.
    if Path(MCDADE_IMPROVE_FILTER_PNG).exists():
        add_image_fit(s, MCDADE_IMPROVE_FILTER_PNG,
                      left=0.5, top=1.35, width=W - 1.0, height=5.0)
        add_takeaway(
            s,
            "IMPROVE filter shows the same array-of-dots deposition pattern from the metal "
            "support screen, same physics behind Warren's pixelation discussion.",
            top=6.55, height=0.85, size=18,
        )
        return

    # Last resort: placeholder card.
    if True:
        # Placeholder card with instructions — easy to swap in the figure later.
        card = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2.0), Inches(2.0), Inches(W - 4.0), Inches(3.5),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xF7, 0xF1, 0xEC)
        card.line.color.rgb = ACCENT
        card.line.width = Pt(1.5)
        # text overlay
        box = s.shapes.add_textbox(
            Inches(2.4), Inches(2.4), Inches(W - 4.8), Inches(2.7)
        )
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = "Placeholder: McDade IMPROVE filter photo"
        for r in tf.paragraphs[0].runs:
            r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = ACCENT
        for line in [
            "",
            "Save the IMPROVE filter pattern image (from the McDade et al. early-2000s",
            "paper Ann referenced) to:",
            "",
            "    tmp_warren_pages/mcdade_improve_filter.png",
            "",
            "and re-run build_deck.py. The slide will auto-populate.",
        ]:
            p = tf.add_paragraph()
            p.text = line
            for r in p.runs:
                r.font.size = Pt(16); r.font.color.rgb = TEXT_DARK
        add_takeaway(
            s,
            "Per Ann (4/30 email): IMPROVE filter pattern from her McDade-first-author paper, "
            "for visual comparison with the SPARTAN photos on the previous slide.",
            top=6.1, height=1.0, size=15,
        )


# ── Generic figure slide ──────────────────────────────────────────────────
def slide_figure(prs, blank, title, png, takeaway, *, fig_top=1.3, fig_h=5.9):
    """Title + figure slide. The takeaway line is written into speaker notes
    rather than rendered on the slide so the projected slide stays clean."""
    s = prs.slides.add_slide(blank)
    add_title(s, title)
    add_accent_underline(s)
    add_image_fit(s, png, left=0.5, top=fig_top, width=W - 1.0, height=fig_h)
    if takeaway:
        add_speaker_notes(s, takeaway)


# ── Slide 6 — Iron: two panels side by side ────────────────────────────────
def slide_iron(prs, blank):
    s = prs.slides.add_slide(blank)
    add_title(s, "Iron: Fe vs Fabs and Fe/EC distribution")
    add_accent_underline(s)
    half = (W - 1.4) / 2
    add_image_fit(s, FIG_DIR / "slide06a_fe_vs_fabs.png",
                  left=0.5, top=1.35, width=half, height=5.7)
    add_image_fit(s, FIG_DIR / "slide06b_fe_ec_ratio.png",
                  left=0.5 + half + 0.4, top=1.35, width=half, height=5.7)
    add_speaker_notes(
        s,
        "Fe absorption can't explain the Addis offset; pixelation is colorless and "
        "Fe/EC at Addis sits inside the IMPROVE 5–95% band.",
    )


# ── Slide 7 — Seasonality + checklist ──────────────────────────────────────
def slide_seasonality(prs, blank):
    s = prs.slides.add_slide(blank)
    add_title(s, "Seasonality and what we have ruled out")
    add_accent_underline(s)
    add_image_fit(s, FIG_DIR / "slide07_seasonality.png",
                  left=0.5, top=1.3, width=W - 1.0, height=5.9)
    add_speaker_notes(
        s,
        "Each site has a seasonal shape: Delhi pre-monsoon, Beijing winter, JPL flat, "
        "Addis steps up dry → belg → kiremt, but none explain the Addis offset alone.\n"
        "Already checked / ruled out: Fe loading and Fe/EC ratio · MA350 smoothing (% difference) · "
        "flow corrections · filter sub-IDs and lot numbers · HIPS MAC = 10 · EC range vs IMPROVE.",
    )


# ── Slide 11 — Questions for Warren ────────────────────────────────────────
def slide_questions(prs, blank):
    s = prs.slides.add_slide(blank)
    add_title(s, "Questions for Warren")
    add_accent_underline(s)
    items = [
        "Does the Addis Fabs offset look like an extreme version of pixelation, or something else entirely?",
        "Pixelation predicts UNDER-estimation at high loadings; in IMPROVE the high-mass-loading "
        "samples have HIGHER Fabs, not lower. How should we read that?",
        "For SPARTAN R/T values that fall outside the IMPROVE envelope, do those filters look "
        "physically plausible to you?",
        "Any unpublished / residual chemical-composition tests (Fe, dust) you would re-run on Addis?",
        "Is there an IMPROVE subset (e.g. wildfire-impacted, very high-load days) we should compare "
        "Addis against directly?",
    ]
    add_bullets(s, items, top=1.5, size=22, line_spacing=1.3)


def build():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    # 1
    slide_intro(prs, blank)
    # 2
    slide_data_tools(prs, blank)
    # 2.5 — Addis concentrations in context (Ann's 5/8 email)
    slide_addis_context(prs, blank)
    # 2a — SPARTAN filter photos (Ann's 4/30 email)
    slide_spartan_filters(prs, blank)
    # 2b — McDade IMPROVE filter pattern (Ann's 4/30 email)
    slide_mcdade_improve(prs, blank)
    # 3
    slide_figure(
        prs, blank,
        "Slide 3:  Fabs (1/Mm) vs FTIR EC concentration (µg/m³)",
        FIG_DIR / "slide03_fabs_vs_ec_concentration.png",
        "Same EC concentration in air, very different Fabs: Addis sits well above the other three sites.",
    )
    # 4
    slide_figure(
        prs, blank,
        "Slide 4:  Fabs (1/Mm) vs FTIR EC mass on filter (µg)",
        FIG_DIR / "slide04_fabs_vs_ftir_ec.png",
        "Same EC mass on the filter (what HIPS actually 'sees') and the Addis offset is still there.",
    )
    # 5
    slide_figure(
        prs, blank,
        "Slide 5:  Fabs (1/Mm) vs Aethalometer IR BCc (µg/m³)",
        FIG_DIR / "slide05_fabs_vs_aethalometer.png",
        "Independent absorption measurement: MA350 broadly tracks FTIR EC. The Addis offset is in HIPS, not in BC.",
    )
    # 6
    slide_iron(prs, blank)
    # 7
    slide_seasonality(prs, blank)
    # 8
    slide_figure(
        prs, blank,
        "Slide 8:  Four SPARTAN sites on IMPROVE (EC mass on filter)",
        FIG_DIR / "slide08_overlay_mass_on_filter.png",
        "High-mass-loading IMPROVE points have HIGHER Fabs, not lower; opposite of what pixelation predicts. "
        "SPARTAN sites layer on top in distinct colors.",
    )
    # 9
    slide_figure(
        prs, blank,
        "Slide 9:  Per site, 5–95 percentile shading on both axes",
        FIG_DIR / "slide09_per_site_shaded.png",
        "Even at the LOW end of Addis's 5–95% EC range, Fabs is still anomalously high. "
        "Addis and IMPROVE essentially do not overlap.",
    )
    # 10 — Scaled HIPS coordinates (t + r = 1 zero-absorption locus), split by site
    slide_figure(
        prs, blank,
        "Slide 10:  Scaled HIPS coordinates  (t + r = 1)",
        FIG_DIR / "slide10b_scaled_hips.png",
        "Field blanks sit on the t + r = 1 line; loaded filters drop below it. "
        "Addis samples drop the furthest, consistent with the high Fabs we see on every other slide.",
    )
    # 10b — All four overlaid
    slide_figure(
        prs, blank,
        "Slide 10b:  All four sites overlaid: Addis transmittance offset",
        FIG_DIR / "slide10d_all_sites_overlay.png",
        "Addis transmittance is systematically below the other three sites at comparable reflectance.",
    )
    # 11
    slide_questions(prs, blank)

    DECK_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(DECK_PATH)
    return DECK_PATH


if __name__ == "__main__":
    p = build()
    print(f"Saved deck: {p}")
