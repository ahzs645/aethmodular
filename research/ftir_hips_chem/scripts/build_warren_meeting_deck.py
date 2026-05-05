"""Assemble the Warren-meeting slide deck from pre-rendered PNG figures.

Run after `warren_meeting_slides.ipynb` has produced figures into
`output/warren_meeting/figures/`. This script consumes those PNGs (plus the
SPARTAN filter photos in ~/Downloads/filter and any reference PDFs in
tmp_warren_pages/) and writes `output/warren_meeting/warren_meeting_slides.pptx`.

Slide order follows Ann Dillner's email feedback (Apr 29 2026) plus her later
note about adding SPARTAN/IMPROVE filter-pattern photos.
"""
from __future__ import annotations

import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

REPO_ROOT = Path(__file__).resolve().parents[3]
RESEARCH_DIR = REPO_ROOT / "research" / "ftir_hips_chem"
OUTPUT_DIR = RESEARCH_DIR / "output" / "warren_meeting"
FIG_DIR = OUTPUT_DIR / "figures"
DECK_PATH = OUTPUT_DIR / "warren_meeting_slides.pptx"

SPARTAN_FILTER_PHOTOS = [
    Path.home() / "Downloads" / "filter" / "image004.jpg",
    Path.home() / "Downloads" / "filter" / "image005.jpg",
    Path.home() / "Downloads" / "filter" / "image006.jpg",
]
WARREN_RT_PNG = REPO_ROOT / "tmp_warren_pages" / "warren_RT.png"
MCDADE_IMPROVE_FILTER_PNG = REPO_ROOT / "tmp_warren_pages" / "mcdade_improve_filter.png"

# Slide dimensions (16:9 widescreen)
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

TITLE_PT = 30
BULLET_PT = 22
CAPTION_PT = 14


def add_title(slide, text: str, *, size: int = TITLE_PT) -> None:
    box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.8))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.bold = True


def add_image(
    slide,
    png_path: Path,
    *,
    left: float = 0.4,
    top: float = 1.1,
    width: float = 12.5,
    height: float = 6.1,
) -> None:
    """Place a PNG inside the (left, top, width, height) box without stretching.

    Preserves the image's native aspect ratio and centres it within the box.
    """
    p = Path(png_path)
    if not p.exists():
        box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = f"[FIGURE MISSING: {png_path}]"
        for run in tf.paragraphs[0].runs:
            run.font.size = Pt(BULLET_PT)
            run.font.italic = True
        return

    with Image.open(p) as im:
        src_w, src_h = im.size
    src_ratio = src_w / src_h
    box_ratio = width / height
    if src_ratio >= box_ratio:
        # Width-bound: fit to box width, scale height accordingly.
        draw_w = width
        draw_h = width / src_ratio
    else:
        # Height-bound: fit to box height, scale width accordingly.
        draw_h = height
        draw_w = height * src_ratio
    draw_left = left + (width - draw_w) / 2
    draw_top = top + (height - draw_h) / 2
    slide.shapes.add_picture(
        str(p),
        Inches(draw_left),
        Inches(draw_top),
        width=Inches(draw_w),
        height=Inches(draw_h),
    )


def add_bullets(
    slide,
    items: list[str],
    *,
    left: float = 0.5,
    top: float = 1.2,
    width: float = 12.3,
    height: float = 5.8,
    size: int = BULLET_PT,
) -> None:
    """Render bullets. A leading-tab item becomes an indented sub-bullet."""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if item.startswith("\t"):
            p.text = "    – " + item.lstrip("\t ")
            p.level = 1
        else:
            p.text = "• " + item
        for run in p.runs:
            run.font.size = Pt(size)


def add_caption(slide, text: str, *, top: float = 7.05) -> None:
    box = slide.shapes.add_textbox(
        Inches(0.4), Inches(top), Inches(12.5), Inches(0.4)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    for run in tf.paragraphs[0].runs:
        run.font.size = Pt(CAPTION_PT)
        run.font.italic = True


def add_footnote_box(
    slide, text: str, *, top: float = 6.55, height: float = 0.85, size: int = 14
) -> None:
    """Bullet-style talking-point block pinned to the bottom of a figure slide."""
    box = slide.shapes.add_textbox(
        Inches(0.4), Inches(top), Inches(12.5), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = "▸ " + text
    for run in tf.paragraphs[0].runs:
        run.font.size = Pt(size)
        run.font.bold = True


def _add_title_slide(prs, blank) -> None:
    """Cover slide with talk title, subtitle, presenter, affiliation, and date."""
    s = prs.slides.add_slide(blank)

    # Vertical accent bar on the left
    from pptx.shapes.autoshape import Shape  # local import keeps top tidy
    from pptx.enum.shapes import MSO_SHAPE
    bar = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.35), Inches(SLIDE_H_IN)
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xF3, 0x9C, 0x12)  # SPARTAN-Addis amber

    # Title
    title_box = s.shapes.add_textbox(
        Inches(0.9), Inches(2.1), Inches(SLIDE_W_IN - 1.5), Inches(1.6)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.text = "Why is HIPS Fabs at Addis Ababa so high?"
    for run in tf.paragraphs[0].runs:
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

    # Subtitle
    sub_box = s.shapes.add_textbox(
        Inches(0.9), Inches(3.7), Inches(SLIDE_W_IN - 1.5), Inches(0.8)
    )
    tf = sub_box.text_frame
    tf.word_wrap = True
    tf.text = "Cross-network comparison with IMPROVE — meeting with Warren White"
    for run in tf.paragraphs[0].runs:
        run.font.size = Pt(22)
        run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # Presenter line
    pres_box = s.shapes.add_textbox(
        Inches(0.9), Inches(5.4), Inches(SLIDE_W_IN - 1.5), Inches(1.4)
    )
    tf = pres_box.text_frame
    tf.word_wrap = True
    tf.text = "Ahmad Jalil"
    for run in tf.paragraphs[0].runs:
        run.font.size = Pt(20)
        run.font.bold = True
    p = tf.add_paragraph()
    p.text = "University of Northern British Columbia · NASA MAIA / SPARTAN"
    for run in p.runs:
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    p = tf.add_paragraph()
    p.text = "May 2026"
    for run in p.runs:
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]

    # ── Slide 0 — title / cover ────────────────────────────────────────────
    _add_title_slide(prs, blank)

    # ── Slide 1 — SPARTAN filter photos (Ann's email) ──────────────────────
    # Three SPARTAN PTFE filter photos arranged across the slide so the
    # support-mesh grid pattern is clearly visible. McDade IMPROVE comparison
    # is appended as a separate slide if the scan is checked in.
    s = prs.slides.add_slide(blank)
    add_title(s, "SPARTAN PTFE filters — support-mesh grid pattern")
    photos_present = [p for p in SPARTAN_FILTER_PHOTOS if p.exists()]
    if photos_present:
        n = len(photos_present)
        gap = 0.3
        total_w = 12.5
        photo_w = (total_w - gap * (n - 1)) / n
        photo_h = 5.4
        photo_top = 1.15
        for i, photo in enumerate(photos_present):
            box_left = 0.4 + i * (photo_w + gap)
            add_image(s, photo, left=box_left, top=photo_top,
                      width=photo_w, height=photo_h)
    add_caption(
        s,
        "SPARTAN filters as collected — the support-mesh grid is the pattern Warren's "
        "pixelation discussion is concerned with. (IMPROVE filter pattern from "
        "McDade et al. early 2000s on next slide if available.)",
    )

    if MCDADE_IMPROVE_FILTER_PNG.exists():
        s = prs.slides.add_slide(blank)
        add_title(s, "IMPROVE filter pattern — McDade et al. (early 2000s)")
        add_image(s, MCDADE_IMPROVE_FILTER_PNG, height=5.6)
        add_caption(
            s,
            "Reference figure from McDade et al. for visual comparison with the SPARTAN "
            "grid on the previous slide.",
        )

    # ── Slide 4 — Fabs vs EC concentration (µg/m³) ─────────────────────────
    s = prs.slides.add_slide(blank)
    add_title(s, "Fabs (1/Mm) vs EC concentration (µg/m³) — four SPARTAN sites")
    add_image(s, FIG_DIR / "slide03_fabs_vs_ec_concentration.png", height=5.6)
    add_footnote_box(
        s,
        "Same EC concentration; HIPS Fabs at Addis sits well above Delhi / JPL / Beijing.",
    )

    # ── Slide 5 — Fabs vs FTIR EC mass on filter ───────────────────────────
    s = prs.slides.add_slide(blank)
    add_title(s, "Fabs (1/Mm) vs FTIR EC mass on filter (µg)")
    add_image(s, FIG_DIR / "slide04_fabs_vs_ftir_ec.png", height=5.6)
    add_footnote_box(
        s,
        "x-axis is mass on filter (µg), so slope is NOT MAC — but the per-filter EC is what HIPS actually 'sees'.",
    )

    # ── Slide 6 — Fabs vs Aethalometer ─────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_title(s, "Fabs (1/Mm) vs Aethalometer IR BCc (µg/m³)")
    add_image(s, FIG_DIR / "slide05_fabs_vs_aethalometer.png", height=5.6)
    add_footnote_box(
        s,
        "Independent absorption measurement; broadly tracks FTIR EC. Adds a third reference for the Addis offset.",
    )

    # ── Slide 7 — Iron: 6a + 6b side-by-side on ONE slide ──────────────────
    s = prs.slides.add_slide(blank)
    add_title(s, "Iron — Fe vs Fabs (left) and Fe/EC distribution (right)")
    half_w = 6.15
    panel_top = 1.15
    panel_h = 5.4
    add_image(
        s, FIG_DIR / "slide06a_fe_vs_fabs.png",
        left=0.35, top=panel_top, width=half_w, height=panel_h,
    )
    add_image(
        s, FIG_DIR / "slide06b_fe_ec_ratio.png",
        left=0.35 + half_w + 0.3, top=panel_top, width=half_w, height=panel_h,
    )
    add_footnote_box(
        s,
        "Pixelation is colorless; Fe should INCREASE Fabs. Fe-range and Fe/EC-range shown so we can compare to IMPROVE Fe roughly.",
    )

    # ── Slide 8 — Seasonality + checklist ──────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_title(s, "Seasonality and what we have ruled out")
    add_image(s, FIG_DIR / "slide07_seasonality.png", top=1.05, height=4.9)
    add_bullets(
        s,
        [
            "Each site has its own seasonal shape — Delhi peaks pre-monsoon, Beijing in winter, "
            "JPL is flat year-round, Addis steps up dry → belg → kiremt.",
            "None of those patterns explain the Addis Fabs offset on its own.",
            "Also ruled out / inconclusive: flow-rate corrections, AE33 smoothing, Fe contamination, "
            "iron / dust loading, sample-week clustering.",
        ],
        top=6.05,
        height=1.35,
        size=14,
    )

    # ── Slide 9 — Four sites overlaid on IMPROVE ───────────────────────────
    s = prs.slides.add_slide(blank)
    add_title(s, "Four SPARTAN sites overlaid on IMPROVE — EC mass on filter")
    add_image(s, FIG_DIR / "slide08_overlay_mass_on_filter.png", height=5.4)
    add_footnote_box(
        s,
        "High-mass-loading IMPROVE samples show HIGHER Fabs, not lower — opposite of what pixelation predicts.",
        top=6.45,
        height=0.95,
        size=15,
    )

    # ── Slide 10 — Per-site, 5–95% shading on both axes ────────────────────
    s = prs.slides.add_slide(blank)
    add_title(s, "Per site — 5–95th percentile shading on both EC mass and Fabs")
    add_image(s, FIG_DIR / "slide09_per_site_shaded.png", height=5.4)
    add_footnote_box(
        s,
        "Even at the low end of Addis's 5–95 percentile EC range, Fabs is anomalously high. "
        "The Addis and IMPROVE distributions essentially do not overlap.",
        top=6.45,
        height=0.95,
        size=15,
    )

    # ── R/T calibration space (3 figures from the notebook; 10a omitted) ───
    s = prs.slides.add_slide(blank)
    add_title(s, "Scaled HIPS coordinates (t + r = 1 zero-absorption locus)")
    add_image(s, FIG_DIR / "slide10b_scaled_hips.png", height=5.6)
    add_footnote_box(
        s,
        "Blanks sit on the t + r = 1 line; loaded filters drop below it. Scale of drop ~ aerosol absorption.",
    )

    s = prs.slides.add_slide(blank)
    add_title(s, "Per site — each filter vs its own row/lot calibration line")
    add_image(s, FIG_DIR / "slide10c_site_panels.png", height=5.6)
    add_footnote_box(
        s,
        "Crimson star = site median. Arrow length = R-suppression vs blank-model R. "
        "Addis: median R-suppression ≈ 62 %, τ ≈ 0.96; way past Delhi/Beijing.",
    )

    s = prs.slides.add_slide(blank)
    add_title(s, "All four sites overlaid — Addis transmittance offset")
    add_image(s, FIG_DIR / "slide10d_all_sites_overlay.png", height=5.6)
    add_footnote_box(
        s,
        "Addis transmittance is systematically below the other three sites at comparable reflectance.",
    )

    if WARREN_RT_PNG.exists():
        s = prs.slides.add_slide(blank)
        add_title(s, "Warren et al. — IMPROVE R / T (for comparison)")
        add_image(s, WARREN_RT_PNG, height=5.6)
        add_caption(
            s,
            "From Warren White et al. (2024). Drop-in if/when IMPROVE R/T export becomes available.",
        )

    # ── Slide 12 — Questions for Warren ────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_title(s, "Questions for Warren")
    add_bullets(
        s,
        [
            "Does the Addis Fabs offset look like an extreme version of pixelation, or something else?",
            "Pixelation predicts UNDER-estimation at high loadings; we see OVER-estimation. How to reconcile?",
            "For SPARTAN R/T values that fall outside the IMPROVE envelope — do those filters look physically plausible?",
            "Any unpublished / residual tests on chemical-composition impacts (Fe, dust) that would be worth re-running on Addis?",
            "Is there a sensible IMPROVE subset (e.g. wildfire-impacted, high-load days) we should compare Addis against?",
        ],
    )

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(DECK_PATH)
    return DECK_PATH


def main() -> int:
    path = build()
    print(f"Saved deck: {path}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
